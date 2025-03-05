import os
import traceback
import fitz
import docx
import pandas as pd
import pptx
import numpy as np
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

# python -c "import nltk; nltk.download('punkt')"
#python -c "import nltk; nltk.download('stopwords')"

class TextProcessor:
    def __init__(self, directory):
        self.directory = directory

    def get_all_file_paths(self):
        """
        Recursively extracts all file paths from the given directory.
        
        :return: A list containing the paths of all files in the directory
        """
        file_paths = []
        for root, _, files in os.walk(self.directory):
            for file in files:
                file_paths.append(os.path.join(root, file))
        return file_paths

    def extract_text_from_file(self, file_path):
        """
        Extract text from various file types.
        
        :param file_path: Path to the file
        :return: Extracted text content or empty string if extraction fails
        """
        if not os.path.exists(file_path):
            print(f"File not found: {file_path}")
            return ""

        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.extract_from_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self.extract_from_word(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self.extract_from_powerpoint(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self.extract_from_excel(file_path)
            elif file_ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']:
                return self.extract_from_text_file(file_path)
            elif file_ext in ['.py', '.js', '.java', '.c', '.cpp', '.h', '.cs', '.php', '.rb', '.go', '.swift']:
                return self.extract_from_text_file(file_path)
        except Exception as e:
            print(f"Error extracting text from {file_path}: {str(e)}")
            print(traceback.format_exc())
            return ""
        
    def clean_text(self, text):
        tokens = word_tokenize(text.lower())
        stop_words = set(stopwords.words('english'))
        filtered_tokens = [token for token in tokens if token.isalnum() and token not in stop_words]
        clean_text = ' '.join(filtered_tokens)
        return clean_text

    def extract_from_pdf(self, file_path):
        try:
            doc = fitz.open(file_path)
            text = "".join([page.get_text() for page in doc[:20]])  # Limit to first 20 pages
            doc.close()
            return text
        except Exception as e:
            print(f"Error extracting from PDF: {str(e)}")
            return ""

    def extract_from_word(self, file_path):
        try:
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"Error extracting from Word: {str(e)}")
            return ""

    def extract_from_powerpoint(self, file_path):
        try:
            prs = pptx.Presentation(file_path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            return text
        except Exception as e:
            print(f"Error extracting from PowerPoint: {str(e)}")
            return ""

    def extract_from_excel(self, file_path):
        try:
            all_dfs = pd.read_excel(file_path, sheet_name=None)
            return "\n".join([f"Sheet: {name}\n{df.to_string(index=False)}" for name, df in all_dfs.items()])
        except Exception as e:
            print(f"Error extracting from Excel: {str(e)}")
            return ""

    def extract_from_text_file(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading text file: {str(e)}")
            return ""